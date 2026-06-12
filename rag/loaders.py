from pathlib import Path


# 把从各种库里提取出来的乱七八糟文本，统一清理成“干净的一块文本”。去掉每行末尾多余空格。保留合理的换行
# 纯 Python 字符串处理，splitlines() + rstrip() 是非常经典的文本清洗技巧。
def _normalize_block(text: str) -> str:
    return "\n".join(line.rstrip() for line in text.splitlines()).strip()


def _load_pdf(filepath: str) -> list[dict]:
    from pypdf import PdfReader

    # PDF 按页切，尽量保留“页码 -> 内容”的对应关系，
    # 后续回答时更容易追溯来源。
    reader = PdfReader(filepath)
    documents = []
    for index, page in enumerate(reader.pages, start=1):
        text = _normalize_block(page.extract_text() or "")
        if not text:
            continue
        documents.append({
            "page_content": text,
            "metadata": {
                "source_file": Path(filepath).name,
                "page": index,
                "slide": None,
                "source_type": "pdf",
            },
        })
    return documents


def _load_docx(filepath: str) -> list[dict]:
    from docx import Document as DocxDocument

    doc = DocxDocument(filepath)
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                # 表格先按行拼成文本，第一版优先保证信息可检索，
                # 暂不保留复杂表格结构。
                parts.append(" | ".join(cells))
    text = _normalize_block("\n".join(parts))
    if not text:
        return []
    return [{
        "page_content": text,
        "metadata": {
            "source_file": Path(filepath).name,
            "page": 1,
            "slide": None,
            "source_type": "docx",
        },
    }]


def _load_pptx(filepath: str) -> list[dict]:
    from pptx import Presentation

    presentation = Presentation(filepath)
    documents = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        # PPT 按页作为逻辑块，避免把一页说明拆散后丢上下文。
        text = _normalize_block("\n".join(filter(None, texts)))
        if not text:
            continue
        documents.append({
            "page_content": text,
            "metadata": {
                "source_file": Path(filepath).name,
                "page": None,
                "slide": index,
                "source_type": "pptx",
            },
        })
    return documents


def _load_doc(filepath: str) -> list[dict]:
    try:
        import pythoncom
        from win32com.client import Dispatch
    except ImportError as exc:
        raise RuntimeError("解析 .doc 需要安装 pywin32，并在 Windows 环境中运行。") from exc

    # 项目资料里实际存在旧版 .doc 文件；标准 Python 生态对其支持较差，
    # 因此这里使用 Windows + Word COM 作为最小可行方案。
    pythoncom.CoInitialize()
    word = Dispatch("Word.Application")
    word.Visible = False
    doc = None
    try:
        doc = word.Documents.Open(str(Path(filepath).resolve()))
        text = _normalize_block(doc.Content.Text or "")
    except Exception as exc:
        raise RuntimeError(f"读取 .doc 失败: {Path(filepath).name}，请确认本机安装了 Microsoft Word。") from exc
    finally:
        if doc is not None:
            doc.Close(False)
        word.Quit()

    if not text:
        return []
    return [{
        "page_content": text,
        "metadata": {
            "source_file": Path(filepath).name,
            "page": 1,
            "slide": None,
            "source_type": "doc",
        },
    }]


def load_document(filepath: str) -> list[dict]:
    path = Path(filepath)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(filepath)
    if suffix == ".docx":
        return _load_docx(filepath)
    if suffix == ".pptx":
        return _load_pptx(filepath)
    if suffix == ".doc":
        return _load_doc(filepath)
    raise ValueError(f"不支持的文档类型: {path.name}")
